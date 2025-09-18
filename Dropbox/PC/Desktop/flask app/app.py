import streamlit as st
import os
import google.generativeai as genai
from dotenv import load_dotenv
import pandas as pd
import PyPDF2
import docx
import pptx
import chromadb
from sentence_transformers import SentenceTransformer
import numpy as np
from pathlib import Path
import json
import tempfile
import time

# Page configuration
st.set_page_config(
    page_title="Semantic Search RAG Application",

    layout="wide",
    initial_sidebar_state="expanded"
)

# Load environment variables
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

# Configure Gemini API
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
else:
    st.error("Gemini API key not found. Please set GEMINI_API_KEY in your .env file.")

# Initialize embedding model
@st.cache_resource
def load_embedding_model():
    return SentenceTransformer('sentence-transformers/all-MiniLM-L6-v2')

# Initialize vector database
@st.cache_resource
def init_vector_db():
    client = chromadb.PersistentClient(path="./database")
    collection = client.get_or_create_collection(
        name="documents",
        metadata={"hnsw:space": "cosine"}
    )
    return collection

# Custom CSS with abstract background and blue-orange theme
st.markdown(
    """
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=Playfair+Display:wght@300;400;500;600;700&display=swap');
    
    * {
        font-family: 'Inter', sans-serif;
    }
    
    html, body, .stApp {
        background: linear-gradient(135deg, #BF360C, #002171, #4A148C, #880E4F, #BF360C, #002171, #4A148C, #880E4F) !important;
        background-size: 300% 300% !important;
        animation: gradientShift 8s ease-in-out infinite !important;
        min-height: 100vh !important;
        overflow-x: hidden !important;
        position: relative !important;
    }
    
    /* Background gradient animation */
    @keyframes gradientShift {
        0% { 
            background-position: 0% 50%;
        }
        25% { 
            background-position: 100% 0%;
        }
        50% { 
            background-position: 200% 50%;
        }
        75% { 
            background-position: 100% 100%;
        }
        100% { 
            background-position: 0% 50%;
        }
    }
    
    .main > div {
        background: transparent !important;
        padding-top: 2rem;
    }
    

    
    .stApp > div {
        background: transparent !important;
    }
    
    .css-1d391kg {
        background: rgba(255, 255, 255, 0.05) !important;
    }
    
    /* Animated gradient border for containers */
    @keyframes borderGlow {
        0% { 
            border-color: #01579B !important;
            box-shadow: 
                0 0 20px rgba(1, 87, 155, 0.6),
                0 0 40px rgba(1, 87, 155, 0.4),
                inset 0 0 20px rgba(1, 87, 155, 0.1) !important;
        }
        25% { 
            border-color: #4A148C !important;
            box-shadow: 
                0 0 25px rgba(74, 20, 140, 0.6),
                0 0 50px rgba(74, 20, 140, 0.4),
                inset 0 0 25px rgba(74, 20, 140, 0.1) !important;
        }
        50% { 
            border-color: #BF360C !important;
            box-shadow: 
                0 0 30px rgba(191, 54, 12, 0.7),
                0 0 60px rgba(191, 54, 12, 0.5),
                inset 0 0 30px rgba(191, 54, 12, 0.1) !important;
        }
        75% { 
            border-color: #4A148C !important;
            box-shadow: 
                0 0 25px rgba(74, 20, 140, 0.6),
                0 0 50px rgba(74, 20, 140, 0.4),
                inset 0 0 25px rgba(74, 20, 140, 0.1) !important;
        }
        100% { 
            border-color: #01579B !important;
            box-shadow: 
                0 0 20px rgba(1, 87, 155, 0.6),
                0 0 40px rgba(1, 87, 155, 0.4),
                inset 0 0 20px rgba(1, 87, 155, 0.1) !important;
        }
    }
    
    @keyframes sidebarGradientShift {
        0% { background: linear-gradient(135deg, rgba(255, 107, 53, 0.08) 0%, rgba(59, 130, 246, 0.08) 100%); }
        50% { background: linear-gradient(135deg, rgba(59, 130, 246, 0.08) 0%, rgba(255, 107, 53, 0.08) 100%); }
        100% { background: linear-gradient(135deg, rgba(255, 107, 53, 0.08) 0%, rgba(59, 130, 246, 0.08) 100%); }
    }
    
    @keyframes blueGlow {
        0% { 
            border: 2px solid #01579B !important;
            box-shadow: 0 0 25px rgba(1, 87, 155, 0.5), 0 0 50px rgba(191, 54, 12, 0.2) !important;
        }
        25% { 
            border: 2px solid #4A148C !important;
            box-shadow: 0 0 30px rgba(74, 20, 140, 0.5), 0 0 60px rgba(1, 87, 155, 0.3) !important;
        }
        50% { 
            border: 2px solid #BF360C !important;
            box-shadow: 0 0 35px rgba(191, 54, 12, 0.6), 0 0 70px rgba(74, 20, 140, 0.3) !important;
        }
        75% { 
            border: 2px solid #4A148C !important;
            box-shadow: 0 0 30px rgba(74, 20, 140, 0.5), 0 0 60px rgba(191, 54, 12, 0.3) !important;
        }
        100% { 
            border: 2px solid #01579B !important;
            box-shadow: 0 0 25px rgba(1, 87, 155, 0.5), 0 0 50px rgba(191, 54, 12, 0.2) !important;
        }
    }
    
    .main-container {
        background: rgba(255, 255, 255, 0.05) !important;
        backdrop-filter: blur(20px) !important;
        border-radius: 20px !important;
        border: 3px solid #01579B !important;
        padding: 3rem !important;
        margin: 2rem 0 !important;
        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.2) !important;
        animation: borderGlow 4s ease-in-out infinite !important;
        min-height: 200px !important;
        position: relative !important;
    }
    
    .search-container {
        background: rgba(255, 255, 255, 0.05) !important;
        backdrop-filter: blur(20px) !important;
        border-radius: 15px !important;
        border: 3px solid #01579B !important;
        padding: 3px !important;
        margin: 1rem 0 !important;
        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.2) !important;
        animation: borderGlow 4s ease-in-out infinite !important;
        position: relative !important;
        width: auto !important;
        max-width: fit-content !important;
        height: 76px !important;
        display: flex !important;
        align-items: center !important;
    }
    
    .stTextInput > div > div > input {
        background: rgba(255, 255, 255, 0.05) !important;
        border: 2px solid #002171 !important;
        border-radius: 12px !important;
        color: white !important;
        font-size: 20px !important;
        font-weight: 300 !important;
        padding: 1.5rem 2rem !important;
        transition: all 0.3s ease !important;
        backdrop-filter: blur(20px) !important;
        height: 64px !important;
        min-height: 64px !important;
        max-height: 64px !important;
        letter-spacing: 0.3px !important;
        box-shadow: 0 0 20px rgba(59, 130, 246, 0.4), 0 8px 32px rgba(0, 0, 0, 0.1) !important;
        animation: blueGlow 3s ease-in-out infinite !important;
        display: flex !important;
        align-items: center !important;
        margin: 0px !important;
        width: 100% !important;
        position: relative !important;
        z-index: 2 !important;
    }
    
    .stTextInput > div > div > input:focus {
        border-color: #3B82F6 !important;
        box-shadow: 0 0 30px rgba(59, 130, 246, 0.6), 0 0 60px rgba(59, 130, 246, 0.3) !important;
        background: rgba(255, 255, 255, 0.08) !important;
    }
    
    .stTextInput > div > div {
        height: 64px !important;
        min-height: 64px !important;
    }
    
    .stTextInput > div {
        height: 64px !important;
        min-height: 64px !important;
    }
    
    .stTextInput > div > div > input::placeholder {
        color: rgba(255, 255, 255, 0.6) !important;
        font-size: 18px !important;
        font-weight: 300 !important;
        letter-spacing: 0.3px !important;
    }
    
    .stButton > button {
        background: rgba(255, 255, 255, 0.05) !important;
        backdrop-filter: blur(20px) !important;
        border: 2px solid rgba(255, 255, 255, 0.1) !important;
        border-radius: 15px !important;
        color: white !important;
        padding: 1.5rem 2rem !important;
        font-weight: 300 !important;
        font-size: 20px !important;
        transition: all 0.3s ease !important;
        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1) !important;
        animation: borderGlow 3s ease-in-out infinite !important;
        position: relative !important;
        overflow: hidden !important;
        height: 70px !important;
        min-height: 70px !important;
        letter-spacing: 0.3px !important;
    }
    
    .stButton > button:before {
        content: '' !important;
        position: absolute !important;
        top: 0 !important;
        left: 0 !important;
        right: 0 !important;
        bottom: 0 !important;
        background: linear-gradient(135deg, #FF6B35 0%, #3B82F6 100%) !important;
        opacity: 0.8 !important;
        z-index: -1 !important;
        border-radius: 13px !important;
    }
    
    .stButton > button:hover {
        transform: translateY(-3px) !important;
        box-shadow: 0 15px 35px rgba(255, 107, 53, 0.6) !important;
        background: rgba(255, 255, 255, 0.1) !important;
    }
    
    .stButton > button:hover:before {
        opacity: 1 !important;
        background: linear-gradient(135deg, #3B82F6 0%, #FF6B35 100%) !important;
    }
    
    .stButton > button:active {
        transform: translateY(-1px) !important;
    }
    
    .response-container {
        background: rgba(255, 255, 255, 0.05);
        backdrop-filter: blur(20px);
        border-radius: 20px;
        border: 2px solid rgba(255, 255, 255, 0.1);
        padding: 2rem;
        margin: 2rem 0;
        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
        animation: borderGlow 3s ease-in-out infinite;
    }
    
    .source-doc {
        background: rgba(255, 255, 255, 0.05);
        backdrop-filter: blur(20px);
        border-radius: 20px;
        border: 2px solid rgba(255, 255, 255, 0.1);
        padding: 2rem;
        margin: 1rem 0;
        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
        animation: borderGlow 3s ease-in-out infinite;
        transition: all 0.3s ease;
    }
    
    .source-doc:hover {
        background: rgba(255, 255, 255, 0.1);
        transform: translateX(10px);
    }

    .sidebar .stFileUploader {
        background: rgba(255, 255, 255, 0.05);
        backdrop-filter: blur(20px);
        border-radius: 20px;
        border: 2px solid rgba(255, 255, 255, 0.1);
        padding: 2rem;
        margin: 1rem 0;
        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
        animation: borderGlow 3s ease-in-out infinite;
    }
    
    /* Animation keyframes */
    @keyframes fadeInUp {
        from {
            opacity: 0;
            transform: translateY(30px);
        }
        to {
            opacity: 1;
            transform: translateY(0);
        }
    }
    
    @keyframes pulse {
        0%, 100% { opacity: 1; }
        50% { opacity: 0.7; }
    }
    
    .stSelectbox > div > div {
        background: rgba(255, 255, 255, 0.1);
        border-radius: 10px;
        border: 1px solid rgba(255, 255, 255, 0.2);
    }
    
    .stExpander {
        background: rgba(255, 255, 255, 0.05);
        border-radius: 15px;
        border: 1px solid rgba(255, 255, 255, 0.1);
        backdrop-filter: blur(10px);
    }
    
    /* Title styling */
    h1 {
        background: linear-gradient(135deg, #FF6B35, #3B82F6);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        background-clip: text;
        text-align: center;
        font-weight: 700;
        font-size: 3rem;
        margin-bottom: 0.5rem;
    }
    
    .elegant-text {
        color: #ffffff !important;
        font-weight: 300 !important;
        font-size: 2.5rem !important;
        text-align: center !important;
        margin: 2rem 0 !important;
        letter-spacing: 0.5px !important;
        line-height: 1.2 !important;
        text-shadow: 0 2px 10px rgba(0, 0, 0, 0.3) !important;
    }
    
    h2, h3 {
        color: #ffffff;
        font-weight: 300;
        letter-spacing: 0.3px;
    }
    
    .stMarkdown p {
        color: rgba(255, 255, 255, 0.9);
        line-height: 1.6;
        font-weight: 300;
        letter-spacing: 0.2px;
    }
    
    /* All text elements styling */
    * {
        font-family: 'Inter', sans-serif;
        font-weight: 300;
        letter-spacing: 0.3px;
    }
    
    /* Sidebar styling */
    .css-1d391kg {
        background: rgba(255, 255, 255, 0.05) !important;
        backdrop-filter: blur(30px) !important;
        border-right: 2px solid #3B82F6 !important;
        box-shadow: 0 0 20px rgba(59, 130, 246, 0.3), 0 0 50px rgba(59, 130, 246, 0.1) !important;
        animation: blueGlow 3s ease-in-out infinite !important;
    }
    
    .css-1d391kg h3 {
        color: rgba(255, 255, 255, 0.95) !important;
        font-family: 'Playfair Display', serif !important;
        font-weight: 500 !important;
        font-size: 1.8rem !important;
        letter-spacing: 0.8px !important;
        text-align: center !important;
        margin-bottom: 2rem !important;
        text-shadow: 0 0 20px rgba(255, 107, 53, 0.3) !important;
    }
    
    .css-1d391kg .stFileUploader {
        background: rgba(255, 255, 255, 0.05) !important;
        backdrop-filter: blur(25px) !important;
        border: 2px solid #3B82F6 !important;
        border-radius: 20px !important;
        padding: 2rem !important;
        margin: 1.5rem 0 !important;
        box-shadow: 0 0 20px rgba(59, 130, 246, 0.3), 0 15px 45px rgba(0, 0, 0, 0.15) !important;
        animation: blueGlow 3s ease-in-out infinite !important;
        transition: all 0.4s ease !important;
    }
    
    .css-1d391kg .stFileUploader:hover {
        transform: translateY(-5px) !important;
        box-shadow: 0 0 30px rgba(59, 130, 246, 0.5), 0 20px 60px rgba(59, 130, 246, 0.2) !important;
        border-color: #60A5FA !important;
    }
    
    .css-1d391kg .stButton > button {
        background: rgba(255, 255, 255, 0.05) !important;
        border: 2px solid #3B82F6 !important;
        border-radius: 15px !important;
        color: white !important;
        font-family: 'Playfair Display', serif !important;
        font-weight: 500 !important;
        font-size: 1.1rem !important;
        letter-spacing: 0.5px !important;
        padding: 1rem 2rem !important;
        transition: all 0.3s ease !important;
        box-shadow: 0 0 20px rgba(59, 130, 246, 0.3) !important;
        animation: blueGlow 3s ease-in-out infinite !important;
    }
    
    .css-1d391kg .stButton > button:hover {
        transform: translateY(-3px) !important;
        box-shadow: 0 0 30px rgba(59, 130, 246, 0.6), 0 15px 40px rgba(59, 130, 246, 0.3) !important;
        background: rgba(255, 255, 255, 0.08) !important;
        border-color: #60A5FA !important;
    }
    
    .uploadedFile {
        background: rgba(255, 255, 255, 0.05);
        backdrop-filter: blur(20px);
        border-radius: 20px;
        border: 2px solid #3B82F6;
        padding: 2rem;
        margin: 1rem 0;
        box-shadow: 0 0 20px rgba(59, 130, 246, 0.3), 0 8px 32px rgba(0, 0, 0, 0.1);
        animation: blueGlow 3s ease-in-out infinite;
    }
    
    /* Progress bar styling */
    .stProgress .st-bo {
        background: linear-gradient(90deg, #FF6B35, #3B82F6);
        border-radius: 10px;
    }
    
    /* Success/Error message styling */
    .stAlert {
        background: rgba(255, 255, 255, 0.05);
        backdrop-filter: blur(20px);
        border-radius: 20px;
        border: 2px solid rgba(255, 255, 255, 0.1);
        padding: 2rem;
        margin: 1rem 0;
        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
        animation: borderGlow 3s ease-in-out infinite;
    }
    </style>
    """,
    unsafe_allow_html=True
)

def extract_text_from_file(uploaded_file):
    """Extract text from uploaded file"""
    text = ""
    
    try:
        if uploaded_file.type == "application/pdf":
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                text += page.extract_text()
        
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = docx.Document(uploaded_file)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        
        elif uploaded_file.type == "text/plain":
            text = str(uploaded_file.read(), "utf-8")
        
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = pptx.Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    
    except Exception as e:
        st.error(f"Error processing file {uploaded_file.name}: {str(e)}")
    
    return text

def chunk_text(text, chunk_size=1000, overlap=200):
    """Split text into overlapping chunks"""
    words = text.split()
    chunks = []
    
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
        
        if i + chunk_size >= len(words):
            break
    
    return chunks

def generate_rag_response(query, context):
    """Generate response using Gemini with retrieved context"""
    prompt = f"""Based on the following context, please answer the question comprehensively and accurately. If the context doesn't contain enough information to answer the question, please say so.

Context:
{context}

Question: {query}

Answer:"""
    
    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"Error generating response: {str(e)}"

def main():
    # Title with beautiful styling
    st.markdown('<h1>Semantic Search RAG Application</h1>', unsafe_allow_html=True)
    st.markdown('<p style="text-align: center; font-size: 1.4rem; color: rgba(255,255,255,0.9); margin-bottom: 2rem; font-family: \'Playfair Display\', serif; font-weight: 400; letter-spacing: 0.5px;">Search smarter & discover faster</p>', unsafe_allow_html=True)
    
    # Initialize components
    embedding_model = load_embedding_model()
    collection = init_vector_db()
    
    # Sidebar for file upload
    with st.sidebar:
        st.markdown('<h3 style="font-family: \'Playfair Display\', serif; font-weight: 500; font-size: 1.8rem; letter-spacing: 0.8px; text-align: center; margin-bottom: 2rem; color: rgba(255,255,255,0.95); text-shadow: 0 0 20px rgba(255, 107, 53, 0.3);">Document Management</h3>', unsafe_allow_html=True)
        
        uploaded_files = st.file_uploader(
            "Upload Documents",
            accept_multiple_files=True,
            type=['pdf', 'docx', 'pptx', 'txt'],
            help="Supported formats: PDF, DOCX, PPTX, TXT"
        )
        
        if uploaded_files:
            st.success(f"{len(uploaded_files)} file(s) selected")
            
            # Show uploaded files
            for file in uploaded_files:
                st.markdown(f"""
                <div class="uploadedFile">
                    {file.name}<br>
                    <small>Size: {file.size / 1024:.1f} KB</small>
                </div>
                """, unsafe_allow_html=True)
            
            if st.button("Process Documents", use_container_width=True):
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                total_files = len(uploaded_files)
                processed_chunks = 0
                
                for idx, uploaded_file in enumerate(uploaded_files):
                    status_text.text(f"Processing {uploaded_file.name}...")
                    
                    # Extract text
                    text = extract_text_from_file(uploaded_file)
                    
                    if text.strip():
                        # Chunk text
                        chunks = chunk_text(text)
                        
                        if chunks:
                            # Generate embeddings
                            embeddings = embedding_model.encode(chunks)
                            
                            # Store in vector database
                            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                                try:
                                    collection.add(
                                        documents=[chunk],
                                        embeddings=[embedding.tolist()],
                                        ids=[f"{uploaded_file.name}_{int(time.time())}_{i}"],
                                        metadatas=[{
                                            "filename": uploaded_file.name,
                                            "chunk_id": i,
                                            "total_chunks": len(chunks),
                                            "file_type": uploaded_file.type
                                        }]
                                    )
                                    processed_chunks += 1
                                except Exception as e:
                                    st.warning(f"Error storing chunk from {uploaded_file.name}: {str(e)}")
                    
                    progress_bar.progress((idx + 1) / total_files)
                
                status_text.text("Processing complete!")
                time.sleep(1)
                status_text.empty()
                progress_bar.empty()
                
                st.success(f"Successfully processed {processed_chunks} document chunks!")
                st.rerun()
        
        # Database management
        st.markdown("---")
        st.markdown('<h3 style="font-family: \'Playfair Display\', serif; font-weight: 500; font-size: 1.6rem; letter-spacing: 0.8px; text-align: center; margin: 2rem 0 1.5rem 0; color: rgba(255,255,255,0.95); text-shadow: 0 0 20px rgba(59, 130, 246, 0.3);">Database Management</h3>', unsafe_allow_html=True)
        
        if st.button("🗑️ Clear Database", use_container_width=True):
            if st.session_state.get('confirm_clear', False):
                try:
                    # Get all IDs and delete them
                    all_items = collection.get()
                    if all_items['ids']:
                        collection.delete(ids=all_items['ids'])
                    st.success("Database cleared!")
                    st.session_state.confirm_clear = False
                    st.rerun()
                except Exception as e:
                    st.error(f"Error clearing database: {str(e)}")
            else:
                st.session_state.confirm_clear = True
                st.warning("Click again to confirm database deletion")
    
    # Main search interface
    st.markdown('<h2 class="elegant-text">Ask Your Documents</h2>', unsafe_allow_html=True)
    
    # Search input with better styling
    col1, col2 = st.columns([5, 1])
    
    with col1:
        query = st.text_input(
            "Search Query",
            placeholder="Ask me anything",
            label_visibility="collapsed"
        )
    
    with col2:
        search_button = st.button("Search", use_container_width=True, type="primary")
    
    # Search functionality
    if search_button and query:
        if not query.strip():
            st.warning("Please enter a search query")
        else:
            try:
                with st.spinner("Searching through your documents..."):
                    # Generate query embedding
                    query_embedding = embedding_model.encode([query])
                    
                    # Search for relevant documents
                    results = collection.query(
                        query_embeddings=query_embedding.tolist(),
                        n_results=5,
                        include=["documents", "metadatas", "distances"]
                    )
                    
                    if results['documents'][0]:
                        # Combine retrieved documents as context
                        context = "\n\n".join(results['documents'][0])
                        
                        # Generate RAG response
                        with st.spinner("🤖 Generating AI response..."):
                            response = generate_rag_response(query, context)
                        
                        # Display results in beautiful containers
                        st.markdown('<div class="response-container">', unsafe_allow_html=True)
                        st.markdown("### 🤖 AI Response")
                        st.markdown(response)
                        st.markdown('</div>', unsafe_allow_html=True)
                        
                        # Show source documents
                        with st.expander("Source Documents", expanded=False):
                            for i, (doc, metadata, distance) in enumerate(zip(
                                results['documents'][0], 
                                results['metadatas'][0], 
                                results['distances'][0]
                            )):
                                similarity = 1 - distance
                                st.markdown(f"""
                                <div class="source-doc">
                                    <h4 style="color: #FF6B35; margin-top: 0;">Source {i+1}</h4>
                                    <p><strong>File:</strong> {metadata.get('filename', 'Unknown')}</p>
                                    <p><strong>Similarity:</strong> {similarity:.1%}</p>
                                    <p><strong>Content:</strong></p>
                                    <p style="font-style: italic; opacity: 0.9;">{doc[:400]}{'...' if len(doc) > 400 else ''}</p>
                                </div>
                                """, unsafe_allow_html=True)
                    else:
                        st.warning("No relevant documents found. Please upload some documents first or try a different search query.")
            
            except Exception as e:
                st.error(f"Search error: {str(e)}")

if __name__ == "__main__":
    main()
